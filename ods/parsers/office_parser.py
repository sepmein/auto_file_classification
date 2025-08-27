"""
Office文档解析器

支持Word、PowerPoint、Excel等Office文档的文本提取
"""

from typing import Dict, Any, Optional
from pathlib import Path
import logging

# Word文档解析
try:
    from docx import Document as DocxDocument

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PowerPoint文档解析
try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Excel文档解析
try:
    import openpyxl

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

# 备用方案：textract
try:
    import textract

    TEXTRACT_AVAILABLE = True
except ImportError:
    TEXTRACT_AVAILABLE = False

from .base_parser import BaseParser, ParsedContent, ParseResult


class OfficeParser(BaseParser):
    """Office文档解析器"""

    def __init__(self, config: Dict[str, Any]):
        super().__init__(config)
        self.supported_extensions = []

        # 根据可用库设置支持的扩展名
        if DOCX_AVAILABLE:
            self.supported_extensions.extend([".docx"])  # python-docx只支持.docx
        if PPTX_AVAILABLE:
            self.supported_extensions.extend([".pptx"])  # python-pptx只支持.pptx
        if OPENPYXL_AVAILABLE:
            self.supported_extensions.extend([".xlsx"])  # openpyxl主要支持.xlsx

        # 如果有textract作为备用，支持更多格式
        if TEXTRACT_AVAILABLE:
            self.supported_extensions.extend(
                [".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx"]
            )

        self.extract_metadata = config.get("office", {}).get("extract_metadata", True)
        self.max_slides = config.get("office", {}).get(
            "max_slides", 50
        )  # PPT最大幻灯片数
        self.max_sheets = config.get("office", {}).get(
            "max_sheets", 10
        )  # Excel最大工作表数

        if not any(
            [DOCX_AVAILABLE, PPTX_AVAILABLE, OPENPYXL_AVAILABLE, TEXTRACT_AVAILABLE]
        ):
            self.logger.error("未安装Office文档解析库，无法解析Office文件")

    def parse(self, file_path: Path) -> ParseResult:
        """
        解析Office文档

        Args:
            file_path: Office文件路径

        Returns:
            ParseResult: 解析结果
        """
        if not self.can_parse(file_path):
            return self.create_error_result(file_path, f"无法解析文件: {file_path}")

        extension = file_path.suffix.lower()

        try:
            # 检查文件大小
            file_size = file_path.stat().st_size
            if file_size == 0:
                return self.create_error_result(file_path, "文件为空")

            # 根据文件类型选择解析方法
            if extension in [".docx", ".doc"]:
                # 检查是否支持.doc格式
                if extension == ".doc" and not TEXTRACT_AVAILABLE:
                    return self.create_error_result(
                        file_path, ".doc格式需要安装textract库: pip install textract"
                    )
                return self._parse_word(file_path)
            elif extension in [".pptx", ".ppt"]:
                # 检查是否支持.ppt格式
                if extension == ".ppt" and not TEXTRACT_AVAILABLE:
                    return self.create_error_result(
                        file_path, ".ppt格式需要安装textract库: pip install textract"
                    )
                return self._parse_powerpoint(file_path)
            elif extension in [".xlsx", ".xls"]:
                # 检查是否支持.xls格式
                if extension == ".xls" and not TEXTRACT_AVAILABLE:
                    return self.create_error_result(
                        file_path, ".xls格式需要安装textract库: pip install textract"
                    )
                return self._parse_excel(file_path)
            else:
                return self.create_error_result(
                    file_path, f"不支持的文件类型: {extension}"
                )

        except Exception as e:
            error_msg = f"Office文档解析失败: {str(e)}"
            self.logger.error(f"{error_msg}, 文件: {file_path}")
            return self.create_error_result(file_path, error_msg)

    def _format_corruption_error_message(
        self, validation_result: Dict[str, str]
    ) -> str:
        """
        格式化文件损坏错误消息

        Args:
            validation_result: 文件验证结果

        Returns:
            str: 格式化的错误消息
        """
        error_msg = f"文件损坏: {validation_result['error']}\n"
        error_msg += f"详细信息: {validation_result['details']}\n\n"
        error_msg += "🔧 恢复建议:\n"

        suggestions = validation_result.get("recovery_suggestions", [])
        for i, suggestion in enumerate(suggestions, 1):
            error_msg += f"{i}. {suggestion}\n"

        error_msg += "\n💡 提示: .docx文件应该是ZIP格式的XML文档，如果文件损坏，可能无法恢复内容。"
        return error_msg

    def _get_unsupported_format_message(self, extension: str) -> str:
        """
        获取不支持格式的详细错误信息

        Args:
            extension: 文件扩展名

        Returns:
            str: 详细的错误信息和安装指导
        """
        format_messages = {
            ".doc": {
                "format": "Microsoft Word 97-2003 (.doc)",
                "library": "textract",
                "install_cmd": "pip install textract",
                "note": "textract 支持多种旧版 Office 格式，包括 .doc, .ppt, .xls",
            },
            ".ppt": {
                "format": "Microsoft PowerPoint 97-2003 (.ppt)",
                "library": "textract",
                "install_cmd": "pip install textract",
                "note": "需要安装 antiword 和 pptx 依赖",
            },
            ".xls": {
                "format": "Microsoft Excel 97-2003 (.xls)",
                "library": "textract",
                "install_cmd": "pip install textract",
                "note": "需要安装 xlrd 依赖",
            },
            ".docm": {
                "format": "Microsoft Word Macro-Enabled (.docm)",
                "library": "python-docx",
                "install_cmd": "pip install python-docx",
                "note": "宏文件可能需要额外的安全考虑",
            },
            ".pptm": {
                "format": "Microsoft PowerPoint Macro-Enabled (.pptm)",
                "library": "python-pptx",
                "install_cmd": "pip install python-pptx",
                "note": "宏文件可能需要额外的安全考虑",
            },
            ".xlsm": {
                "format": "Microsoft Excel Macro-Enabled (.xlsm)",
                "library": "openpyxl",
                "install_cmd": "pip install openpyxl",
                "note": "宏文件可能需要额外的安全考虑",
            },
        }

        if extension in format_messages:
            info = format_messages[extension]
            return (
                f"不支持的文件格式: {info['format']}\n"
                f"需要安装 {info['library']} 库: {info['install_cmd']}\n"
                f"注意: {info['note']}\n"
                f"安装后请重启应用程序。"
            )
        else:
            return (
                f"不支持的文件格式: {extension}\n"
                f"当前支持的 Office 格式: .docx, .pptx, .xlsx\n"
                f"对于旧版格式(.doc, .ppt, .xls)，请安装 textract 库:\n"
                f"pip install textract"
            )

    def get_supported_formats_info(self) -> Dict[str, Any]:
        """
        获取支持格式的详细信息

        Returns:
            Dict[str, Any]: 支持格式信息
        """
        supported_formats = {
            "native_supported": [],
            "textract_supported": [],
            "not_supported": [],
        }

        # 检查原生支持的格式
        if DOCX_AVAILABLE:
            supported_formats["native_supported"].append(
                {
                    "extension": ".docx",
                    "format": "Microsoft Word (.docx)",
                    "library": "python-docx",
                }
            )

        if PPTX_AVAILABLE:
            supported_formats["native_supported"].append(
                {
                    "extension": ".pptx",
                    "format": "Microsoft PowerPoint (.pptx)",
                    "library": "python-pptx",
                }
            )

        if OPENPYXL_AVAILABLE:
            supported_formats["native_supported"].append(
                {
                    "extension": ".xlsx",
                    "format": "Microsoft Excel (.xlsx)",
                    "library": "openpyxl",
                }
            )

        # 检查textract支持的格式
        if TEXTRACT_AVAILABLE:
            supported_formats["textract_supported"] = [
                {"extension": ".doc", "format": "Microsoft Word 97-2003"},
                {"extension": ".ppt", "format": "Microsoft PowerPoint 97-2003"},
                {"extension": ".xls", "format": "Microsoft Excel 97-2003"},
            ]

        # 不支持的格式
        supported_formats["not_supported"] = [
            {"extension": ".docm", "reason": "需要 python-docx"},
            {"extension": ".pptm", "reason": "需要 python-pptx"},
            {"extension": ".xlsm", "reason": "需要 openpyxl"},
            {"extension": ".rtf", "reason": "需要额外解析器"},
        ]

        return supported_formats

    def check_file_integrity(self, file_path: Path) -> Dict[str, Any]:
        """
        检查文件完整性并提供详细报告

        Args:
            file_path: 文件路径

        Returns:
            Dict[str, Any]: 文件完整性报告
        """
        report = {
            "file_path": str(file_path),
            "file_name": file_path.name,
            "extension": file_path.suffix.lower(),
            "exists": file_path.exists(),
            "readable": False,
            "size_bytes": 0,
            "integrity_status": "unknown",
            "issues": [],
            "recommendations": [],
        }

        if not file_path.exists():
            report["integrity_status"] = "missing"
            report["issues"].append("文件不存在")
            report["recommendations"].append("检查文件路径是否正确")
            return report

        try:
            # 基本文件信息
            stat = file_path.stat()
            report["size_bytes"] = stat.st_size
            report["readable"] = True

            # 大小检查
            if stat.st_size == 0:
                report["integrity_status"] = "empty"
                report["issues"].append("文件为空")
                report["recommendations"].append("文件可能已损坏或未正确保存")
                return report

            if stat.st_size < 100:  # 过于小的文件
                report["integrity_status"] = "suspicious"
                report["issues"].append("文件过小")
                report["recommendations"].append("检查文件是否完整")

            # 格式特定检查
            if file_path.suffix.lower() == ".docx":
                validation = self._validate_docx_file(file_path)
                if validation["valid"]:
                    report["integrity_status"] = "good"
                else:
                    report["integrity_status"] = "corrupted"
                    report["issues"].append(validation["error"])
                    report["recommendations"].extend(
                        validation.get("recovery_suggestions", [])
                    )

            else:
                # 对于其他格式，基本检查通过
                report["integrity_status"] = "good"

        except Exception as e:
            report["integrity_status"] = "error"
            report["issues"].append(f"完整性检查失败: {str(e)}")
            report["recommendations"].append("文件可能已损坏或权限不足")

        return report

    def _validate_docx_file(self, file_path: Path) -> Dict[str, str]:
        """
        全面验证.docx文件格式

        Returns:
            Dict[str, str]: {"valid": bool, "error": error_message, "details": detailed_info}
        """
        import zipfile

        try:
            # 检查文件基本属性
            if not file_path.exists():
                return {
                    "valid": False,
                    "error": "文件不存在",
                    "details": str(file_path),
                }

            file_size = file_path.stat().st_size
            if file_size == 0:
                return {
                    "valid": False,
                    "error": "文件为空",
                    "details": f"大小: {file_size} bytes",
                }

            if file_size < 1000:  # .docx文件最小合理大小
                return {
                    "valid": False,
                    "error": "文件过小",
                    "details": f"大小: {file_size} bytes（可能不是有效的.docx文件）",
                }

            # 尝试打开为ZIP文件
            try:
                with zipfile.ZipFile(str(file_path), "r") as zip_ref:
                    # 检查必要的.docx结构文件
                    required_files = ["word/document.xml", "[Content_Types].xml"]

                    missing_files = []
                    for required_file in required_files:
                        if required_file not in zip_ref.namelist():
                            missing_files.append(required_file)

                    if missing_files:
                        return {
                            "valid": False,
                            "error": "无效的.docx文件结构",
                            "details": f"缺少必要文件: {', '.join(missing_files)}",
                        }

                    # 检查word目录
                    word_files = [
                        f for f in zip_ref.namelist() if f.startswith("word/")
                    ]
                    if not word_files:
                        return {
                            "valid": False,
                            "error": "无效的.docx文件结构",
                            "details": "缺少word目录或内容",
                        }

                    return {
                        "valid": True,
                        "error": "",
                        "details": f"文件大小: {file_size} bytes, 包含 {len(zip_ref.namelist())} 个文件",
                    }

            except zipfile.BadZipFile as e:
                return {
                    "valid": False,
                    "error": ".docx文件损坏或不是有效的ZIP文件",
                    "details": f"ZIP错误: {str(e)}",
                    "recovery_suggestions": [
                        "文件可能在传输过程中损坏",
                        "尝试重新下载或从备份恢复文件",
                        "在Office中打开文件并另存为新文件",
                        "检查文件是否被其他程序锁定或修改",
                    ],
                }

        except Exception as e:
            return {
                "valid": False,
                "error": "文件验证失败",
                "details": f"意外错误: {str(e)}",
            }

    def _parse_word(self, file_path: Path) -> ParseResult:
        """
        解析Word文档

        Args:
            file_path: Word文件路径

        Returns:
            ParseResult: 解析结果
        """
        text_parts = []
        metadata = {}

        try:
            # 增强的文件验证
            if file_path.suffix.lower() == ".docx":
                validation = self._validate_docx_file(file_path)
                if not validation["valid"]:
                    error_msg = self._format_corruption_error_message(validation)
                    return self.create_error_result(file_path, error_msg)

            if file_path.suffix.lower() == ".docx" and DOCX_AVAILABLE:
                # 使用python-docx解析.docx文件
                try:
                    doc = DocxDocument(str(file_path))

                    # 验证文档是否成功加载
                    if not hasattr(doc, "paragraphs"):
                        raise Exception("文档对象无效，缺少paragraphs属性")

                except Exception as docx_error:
                    error_str = str(docx_error)
                    if "Package not found" in error_str:
                        return self.create_error_result(
                            file_path,
                            ".docx文件可能已损坏、密码保护或格式不兼容。请尝试: 1) 重新保存文件 2) 检查文件是否损坏 3) 确认文件不是密码保护的",
                        )
                    elif "BadZipFile" in error_str:
                        return self.create_error_result(
                            file_path,
                            ".docx文件不是有效的ZIP格式，文件可能已损坏。请尝试重新下载或保存文件。",
                        )
                    else:
                        return self.create_error_result(
                            file_path, f"Word文档解析失败: {error_str}"
                        )

                # 提取文本
                paragraph_count = 0
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_parts.append(paragraph.text)
                        paragraph_count += 1

                # 提取表格内容
                table_count = 0
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text)
                        if row_text:
                            text_parts.append(" | ".join(row_text))
                            table_count += 1

                # 提取元数据
                if self.extract_metadata and hasattr(doc, "core_properties"):
                    props = doc.core_properties
                    metadata.update(
                        {
                            "title": props.title,
                            "author": props.author,
                            "subject": props.subject,
                            "keywords": props.keywords,
                            "creation_date": (
                                str(props.created) if props.created else None
                            ),
                            "modification_date": (
                                str(props.modified) if props.modified else None
                            ),
                            "last_modified_by": props.last_modified_by,
                            "revision": props.revision,
                            "paragraph_count": paragraph_count,
                            "table_count": table_count,
                        }
                    )

            elif file_path.suffix.lower() == ".doc" and TEXTRACT_AVAILABLE:
                # 使用textract解析.doc文件
                text = textract.process(str(file_path)).decode("utf-8")
                text_parts = [text]

            elif TEXTRACT_AVAILABLE:
                # 使用textract作为备用方案
                text = textract.process(str(file_path)).decode("utf-8")
                text_parts = [text]

            else:
                error_msg = self._get_unsupported_format_message(
                    file_path.suffix.lower()
                )
                return self.create_error_result(file_path, error_msg)

            if not text_parts:
                return self.create_error_result(file_path, "Word文档为空或无法提取文本")

            text = "\n".join(text_parts)
            text = self.clean_text(text)

            content = ParsedContent(
                text=text,
                title=self._extract_title_from_metadata_or_text(
                    metadata, text, file_path.name
                ),
                author=metadata.get("author"),
                creation_date=metadata.get("creation_date"),
                modification_date=metadata.get("modification_date"),
                metadata={**self.get_file_metadata(file_path), **metadata},
            )

            return self.create_success_result(file_path, content)

        except Exception as e:
            raise Exception(f"Word文档解析失败: {e}")

    def _parse_powerpoint(self, file_path: Path) -> ParseResult:
        """
        解析PowerPoint文档

        Args:
            file_path: PowerPoint文件路径

        Returns:
            ParseResult: 解析结果
        """
        text_parts = []
        metadata = {}

        try:
            if file_path.suffix.lower() == ".pptx" and PPTX_AVAILABLE:
                # 使用python-pptx解析.pptx文件
                prs = Presentation(str(file_path))

                # 提取文本
                slide_count = 0
                for slide in prs.slides:
                    if slide_count >= self.max_slides:
                        break

                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text)

                    if slide_texts:
                        text_parts.append(
                            f"幻灯片 {slide_count + 1}:\n" + "\n".join(slide_texts)
                        )

                    slide_count += 1

                # 提取元数据
                if self.extract_metadata and hasattr(prs, "core_properties"):
                    props = prs.core_properties
                    metadata.update(
                        {
                            "title": props.title,
                            "author": props.author,
                            "subject": props.subject,
                            "keywords": props.keywords,
                            "creation_date": (
                                str(props.created) if props.created else None
                            ),
                            "modification_date": (
                                str(props.modified) if props.modified else None
                            ),
                            "last_modified_by": props.last_modified_by,
                            "slide_count": len(prs.slides),
                        }
                    )

            elif TEXTRACT_AVAILABLE:
                # 使用textract作为备用方案
                text = textract.process(str(file_path)).decode("utf-8")
                text_parts = [text]

            else:
                raise Exception("没有可用的PowerPoint解析库")

            if not text_parts:
                return self.create_error_result(
                    file_path, "PowerPoint文档为空或无法提取文本"
                )

            text = "\n\n".join(text_parts)
            text = self.clean_text(text)

            content = ParsedContent(
                text=text,
                title=self._extract_title_from_metadata_or_text(
                    metadata, text, file_path.name
                ),
                author=metadata.get("author"),
                creation_date=metadata.get("creation_date"),
                modification_date=metadata.get("modification_date"),
                page_count=metadata.get("slide_count"),
                metadata={**self.get_file_metadata(file_path), **metadata},
            )

            return self.create_success_result(file_path, content)

        except Exception as e:
            raise Exception(f"PowerPoint文档解析失败: {e}")

    def _parse_excel(self, file_path: Path) -> ParseResult:
        """
        解析Excel文档

        Args:
            file_path: Excel文件路径

        Returns:
            ParseResult: 解析结果
        """
        text_parts = []
        metadata = {}

        try:
            if file_path.suffix.lower() == ".xlsx" and OPENPYXL_AVAILABLE:
                # 使用openpyxl解析.xlsx文件
                workbook = openpyxl.load_workbook(str(file_path), read_only=True)

                sheet_count = 0
                for sheet_name in workbook.sheetnames:
                    if sheet_count >= self.max_sheets:
                        break

                    worksheet = workbook[sheet_name]
                    sheet_texts = [f"工作表: {sheet_name}"]

                    # 提取前几行作为表头和数据样本
                    row_count = 0
                    for row in worksheet.iter_rows(values_only=True):
                        if row_count >= 20:  # 限制每个工作表最多读取20行
                            break

                        # 过滤空值，组合成文本
                        row_values = [str(cell) for cell in row if cell is not None]
                        if row_values:
                            sheet_texts.append(" | ".join(row_values))

                        row_count += 1

                    if len(sheet_texts) > 1:  # 除了表名外还有内容
                        text_parts.append("\n".join(sheet_texts))

                    sheet_count += 1

                workbook.close()

                # Excel元数据相对简单
                metadata.update(
                    {
                        "sheet_count": len(workbook.sheetnames),
                        "sheet_names": workbook.sheetnames[: self.max_sheets],
                    }
                )

            elif TEXTRACT_AVAILABLE:
                # 使用textract作为备用方案
                text = textract.process(str(file_path)).decode("utf-8")
                text_parts = [text]

            else:
                raise Exception("没有可用的Excel解析库")

            if not text_parts:
                return self.create_error_result(
                    file_path, "Excel文档为空或无法提取文本"
                )

            text = "\n\n".join(text_parts)
            text = self.clean_text(text)

            content = ParsedContent(
                text=text,
                title=self._extract_title_from_metadata_or_text(
                    metadata, text, file_path.name
                ),
                creation_date=metadata.get("creation_date"),
                modification_date=metadata.get("modification_date"),
                metadata={**self.get_file_metadata(file_path), **metadata},
            )

            return self.create_success_result(file_path, content)

        except Exception as e:
            raise Exception(f"Excel文档解析失败: {e}")

    def _extract_title_from_metadata_or_text(
        self, metadata: Dict[str, Any], text: str, file_name: str
    ) -> Optional[str]:
        """
        从元数据或文本中提取标题

        Args:
            metadata: 文档元数据
            text: 文档文本
            file_name: 文件名

        Returns:
            Optional[str]: 提取的标题
        """
        # 优先使用元数据中的标题
        if metadata.get("title"):
            title = metadata["title"].strip()
            if title and title.lower() not in ["document", "presentation", "workbook"]:
                return title

        # 其次从文本内容提取
        text_title = self.extract_title_from_text(text, file_name)
        if text_title:
            return text_title

        # 最后使用文件名
        return Path(file_name).stem

    def can_parse(self, file_path: Path) -> bool:
        """
        检查是否可以解析Office文件

        Args:
            file_path: 文件路径

        Returns:
            bool: 是否可以解析
        """
        if not super().can_parse(file_path):
            return False

        extension = file_path.suffix.lower()

        # 检查特定格式的支持
        if extension in [".docx"] and not DOCX_AVAILABLE:
            return TEXTRACT_AVAILABLE
        elif extension in [".pptx"] and not PPTX_AVAILABLE:
            return TEXTRACT_AVAILABLE
        elif extension in [".xlsx"] and not OPENPYXL_AVAILABLE:
            return TEXTRACT_AVAILABLE
        elif extension in [".doc", ".ppt", ".xls"]:
            return TEXTRACT_AVAILABLE

        return True
